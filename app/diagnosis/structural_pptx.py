"""Structural analysis for PPTX files using python-pptx."""

from __future__ import annotations

import asyncio

from app.core.config import settings
from app.core.log import logger
from app.schema.diagnosis import (
    DiagnosisIssue,
    IssueSeverity,
    IssueSource,
    IssueType,
)

__all__ = ("analyze_pptx",)

_EMU_PER_INCH = 914400

_SAFE_FONTS = {
    "Arial", "Calibri", "Cambria", "Courier New", "Georgia",
    "Helvetica", "Tahoma", "Times New Roman", "Verdana",
    "Consolas", "Segoe UI", "Trebuchet MS",
}

# Common print-friendly slide sizes (width, height in EMU)
# A4 landscape: 10" x 7.5"
# Letter landscape: 10" x 7.5"
# Standard (4:3): 10" x 7.5"
_PRINT_FRIENDLY_RATIOS = {
    (4, 3),   # Standard
    (3, 4),   # Standard portrait
}


async def analyze_pptx(file_path: str, job_id: str) -> list[DiagnosisIssue]:
    """Perform structural analysis on a PPTX file."""
    logger.info(f"Job {job_id}: running PPTX structural analysis on {file_path}")
    return await asyncio.to_thread(_analyze_pptx_sync, file_path, job_id)


def _analyze_pptx_sync(file_path: str, job_id: str) -> list[DiagnosisIssue]:
    from pptx import Presentation

    issues: list[DiagnosisIssue] = []
    try:
        prs = Presentation(file_path)
        issues.extend(_check_slide_size(prs))
        issues.extend(_check_fonts(prs))
        issues.extend(_check_margins(prs))
        issues.extend(_check_text_overflow(prs))
    except Exception:
        logger.exception(f"Job {job_id}: PPTX structural analysis failed")
    return issues


def _check_slide_size(prs) -> list[DiagnosisIssue]:
    """Check if slide dimensions are print-friendly."""
    issues: list[DiagnosisIssue] = []

    width_emu = prs.slide_width
    height_emu = prs.slide_height

    if not width_emu or not height_emu:
        return issues

    width_in = width_emu / _EMU_PER_INCH
    height_in = height_emu / _EMU_PER_INCH

    # 16:9 widescreen is not great for print (too wide for most paper)
    ratio = width_emu / height_emu
    if abs(ratio - 16 / 9) < 0.05:
        issues.append(DiagnosisIssue(
            type=IssueType.slide_size_mismatch,
            severity=IssueSeverity.warning,
            source=IssueSource.structural,
            description=(
                f"Slide size is 16:9 widescreen ({width_in:.1f}\" x {height_in:.1f}\") — "
                f"may not fit well on standard paper sizes"
            ),
            suggested_fix="set_pptx_slide_size",
            confidence=0.8,
        ))

    return issues


def _check_fonts(prs) -> list[DiagnosisIssue]:
    """Check for uncommon fonts and small font sizes."""
    issues: list[DiagnosisIssue] = []
    min_pt = settings.DIAGNOSIS_MIN_FONT_PT
    fonts_used: set[str] = set()
    small_font_locations: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts_used.add(run.font.name)
                    if run.font.size is not None:
                        pt = run.font.size / 12700  # EMU to pt
                        if pt < min_pt and run.text.strip():
                            small_font_locations.append(
                                f"slide {slide_num} ({pt:.1f}pt)"
                            )

    uncommon = fonts_used - _SAFE_FONTS
    for font in uncommon:
        issues.append(DiagnosisIssue(
            type=IssueType.non_embedded_font,
            severity=IssueSeverity.warning,
            source=IssueSource.structural,
            description=f"Font '{font}' may not be available on print server",
            suggested_fix="replace_font",
            confidence=0.6,
        ))

    if small_font_locations:
        sample = small_font_locations[:5]
        issues.append(DiagnosisIssue(
            type=IssueType.small_font,
            severity=IssueSeverity.warning,
            source=IssueSource.structural,
            description=(
                f"Small font sizes detected (< {min_pt}pt) in: "
                + ", ".join(sample)
                + (f" and {len(small_font_locations) - 5} more"
                   if len(small_font_locations) > 5 else "")
            ),
            suggested_fix="adjust_pptx_font_size",
            confidence=0.85,
        ))

    return issues


def _check_margins(prs) -> list[DiagnosisIssue]:
    """Check if shapes are placed too close to slide edges."""
    issues: list[DiagnosisIssue] = []

    slide_w = prs.slide_width or 0
    slide_h = prs.slide_height or 0
    if not slide_w or not slide_h:
        return issues

    # Minimum margin from edge (0.25 inch in EMU)
    min_margin_emu = int(0.25 * _EMU_PER_INCH)
    edge_shapes: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            left = shape.left or 0
            top = shape.top or 0
            right = left + (shape.width or 0)
            bottom = top + (shape.height or 0)

            near_edge = (
                left < min_margin_emu
                or top < min_margin_emu
                or right > slide_w - min_margin_emu
                or bottom > slide_h - min_margin_emu
            )
            if near_edge:
                edge_shapes.append(f"slide {slide_num}: {shape.shape_type}")

    if edge_shapes:
        sample = edge_shapes[:5]
        issues.append(DiagnosisIssue(
            type=IssueType.text_outside_printable,
            severity=IssueSeverity.warning,
            source=IssueSource.structural,
            description=(
                f"{len(edge_shapes)} shape(s) placed within 0.25\" of slide edge: "
                + ", ".join(sample)
                + (f" and {len(edge_shapes) - 5} more"
                   if len(edge_shapes) > 5 else "")
            ),
            confidence=0.7,
        ))

    return issues


def _check_text_overflow(prs) -> list[DiagnosisIssue]:
    """Check for text frames that may overflow their bounding boxes."""
    issues: list[DiagnosisIssue] = []
    overflow_count = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame

            # Check if auto-size is set to none (fixed size) with wrapping off
            if tf.word_wrap is False:
                text_len = sum(len(p.text) for p in tf.paragraphs)
                if text_len > 50:  # Arbitrary but reasonable threshold
                    overflow_count += 1

    if overflow_count:
        issues.append(DiagnosisIssue(
            type=IssueType.text_overflow,
            severity=IssueSeverity.warning,
            source=IssueSource.structural,
            description=(
                f"{overflow_count} text frame(s) have word wrap disabled with "
                f"long text — content may overflow when printed"
            ),
            suggested_fix="adjust_pptx_font_size",
            confidence=0.65,
        ))

    return issues
