"""PPTX fix tools: slide size and font adjustments."""

from __future__ import annotations

import asyncio

from app.core.log import logger
from app.schema.fix import FixResult

__all__ = (
    "adjust_pptx_font_size",
    "replace_pptx_font",
    "reposition_pptx_shapes",
    "resize_pptx_text_boxes",
    "set_pptx_slide_size",
)

_EMU_PER_INCH = 914400


async def set_pptx_slide_size(
    file_path: str,
    job_id: str,
    width: float = 10.0,
    height: float = 7.5,
) -> FixResult:
    """Set slide dimensions. Values in inches. Default is 4:3 standard (10x7.5)."""
    return await asyncio.to_thread(
        _set_pptx_slide_size_sync, file_path, job_id, width, height,
    )


def _set_pptx_slide_size_sync(
    file_path: str,
    job_id: str,
    width: float,
    height: float,
) -> FixResult:
    from pptx import Presentation

    try:
        prs = Presentation(file_path)

        old_w = prs.slide_width / _EMU_PER_INCH if prs.slide_width else 0
        old_h = prs.slide_height / _EMU_PER_INCH if prs.slide_height else 0

        prs.slide_width = int(width * _EMU_PER_INCH)
        prs.slide_height = int(height * _EMU_PER_INCH)

        prs.save(file_path)

        return FixResult(
            tool_name="set_pptx_slide_size",
            job_id=job_id,
            success=True,
            description=(
                f"Changed slide size from {old_w:.1f}\"x{old_h:.1f}\" "
                f"to {width:.1f}\"x{height:.1f}\""
            ),
            before_value=f"{old_w:.1f}x{old_h:.1f}",
            after_value=f"{width:.1f}x{height:.1f}",
        )
    except Exception as exc:
        logger.error(f"Job {job_id}: set_pptx_slide_size failed: {exc}")
        return FixResult(
            tool_name="set_pptx_slide_size",
            job_id=job_id,
            success=False,
            description=f"Failed to set PPTX slide size: {exc}",
            error=str(exc),
        )


async def adjust_pptx_font_size(
    file_path: str,
    job_id: str,
    min_size_pt: float = 10.0,
) -> FixResult:
    """Enforce a minimum font size across all text in the presentation."""
    return await asyncio.to_thread(
        _adjust_pptx_font_size_sync, file_path, job_id, min_size_pt,
    )


def _adjust_pptx_font_size_sync(
    file_path: str,
    job_id: str,
    min_size_pt: float,
) -> FixResult:
    from pptx import Presentation
    from pptx.util import Pt

    try:
        prs = Presentation(file_path)
        adjusted = 0
        min_emu = int(min_size_pt * 12700)

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size is not None and run.font.size < min_emu:
                            run.font.size = Pt(min_size_pt)
                            adjusted += 1

        prs.save(file_path)

        if adjusted == 0:
            return FixResult(
                tool_name="adjust_pptx_font_size",
                job_id=job_id,
                success=True,
                description=f"No text runs below {min_size_pt}pt found",
            )

        return FixResult(
            tool_name="adjust_pptx_font_size",
            job_id=job_id,
            success=True,
            description=f"Adjusted {adjusted} text run(s) to minimum {min_size_pt}pt",
            after_value=f"min {min_size_pt}pt",
        )
    except Exception as exc:
        logger.error(f"Job {job_id}: adjust_pptx_font_size failed: {exc}")
        return FixResult(
            tool_name="adjust_pptx_font_size",
            job_id=job_id,
            success=False,
            description=f"Failed to adjust PPTX font size: {exc}",
            error=str(exc),
        )


# ── Shape repositioning ─────────────────────────────────────────────────


async def reposition_pptx_shapes(
    file_path: str,
    job_id: str,
    margin_inches: float = 0.25,
) -> FixResult:
    """Move shapes that extend beyond the printable area back inside slide bounds.

    Shifts shapes inward while preserving their size and relative positions.
    margin_inches defines the safe margin from slide edges.
    """
    return await asyncio.to_thread(
        _reposition_pptx_shapes_sync, file_path, job_id, margin_inches,
    )


def _reposition_pptx_shapes_sync(
    file_path: str,
    job_id: str,
    margin_inches: float,
) -> FixResult:
    from pptx import Presentation

    try:
        prs = Presentation(file_path)
        margin_emu = int(margin_inches * _EMU_PER_INCH)
        slide_w = prs.slide_width or int(10 * _EMU_PER_INCH)
        slide_h = prs.slide_height or int(7.5 * _EMU_PER_INCH)
        repositioned = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                left = shape.left or 0
                top = shape.top or 0
                width = shape.width or 0
                height = shape.height or 0
                moved = False

                # Check and fix left overflow
                if left < margin_emu:
                    shape.left = margin_emu
                    moved = True
                elif left + width > slide_w - margin_emu:
                    # Shift left so right edge is within bounds
                    new_left = slide_w - margin_emu - width
                    if new_left < margin_emu:
                        # Shape is wider than printable area — clamp to left margin
                        new_left = margin_emu
                    shape.left = new_left
                    moved = True

                # Check and fix top overflow
                if top < margin_emu:
                    shape.top = margin_emu
                    moved = True
                elif top + height > slide_h - margin_emu:
                    new_top = slide_h - margin_emu - height
                    if new_top < margin_emu:
                        new_top = margin_emu
                    shape.top = new_top
                    moved = True

                if moved:
                    repositioned += 1

        prs.save(file_path)

        return FixResult(
            tool_name="reposition_pptx_shapes",
            job_id=job_id,
            success=True,
            description=(
                f"Repositioned {repositioned} shape(s) within "
                f'{margin_inches}" margin'
            ),
            after_value=f"{repositioned} shapes moved",
        )
    except Exception as exc:
        logger.error(f"Job {job_id}: reposition_pptx_shapes failed: {exc}")
        return FixResult(
            tool_name="reposition_pptx_shapes",
            job_id=job_id,
            success=False,
            description=f"Failed to reposition PPTX shapes: {exc}",
            error=str(exc),
        )


# ── Font replacement ─────────────────────────────────────────────────────


async def replace_pptx_font(
    file_path: str,
    job_id: str,
    from_font: str,
    to_font: str,
) -> FixResult:
    """Replace all occurrences of a font across all text in the presentation."""
    return await asyncio.to_thread(
        _replace_pptx_font_sync, file_path, job_id, from_font, to_font,
    )


def _replace_pptx_font_sync(
    file_path: str,
    job_id: str,
    from_font: str,
    to_font: str,
) -> FixResult:
    from pptx import Presentation

    try:
        prs = Presentation(file_path)
        replaced = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name == from_font:
                            run.font.name = to_font
                            replaced += 1

                # Also check table cells if present
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    if run.font.name == from_font:
                                        run.font.name = to_font
                                        replaced += 1

        prs.save(file_path)

        return FixResult(
            tool_name="replace_pptx_font",
            job_id=job_id,
            success=True,
            description=f"Replaced '{from_font}' with '{to_font}' in {replaced} run(s)",
            before_value=from_font,
            after_value=to_font,
        )
    except Exception as exc:
        logger.error(f"Job {job_id}: replace_pptx_font failed: {exc}")
        return FixResult(
            tool_name="replace_pptx_font",
            job_id=job_id,
            success=False,
            description=f"Failed to replace PPTX font: {exc}",
            error=str(exc),
        )


# ── Text box resizing ────────────────────────────────────────────────────


async def resize_pptx_text_boxes(
    file_path: str,
    job_id: str,
    strategy: str = "shrink_text",
) -> FixResult:
    """Resize text boxes to fit their content.

    Strategies:
        - 'shrink_text': Enable auto-size (shrink text to fit box)
        - 'grow': Expand text box dimensions to fit content
    """
    return await asyncio.to_thread(
        _resize_pptx_text_boxes_sync, file_path, job_id, strategy,
    )


def _resize_pptx_text_boxes_sync(
    file_path: str,
    job_id: str,
    strategy: str,
) -> FixResult:
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE

    try:
        prs = Presentation(file_path)
        adjusted = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                tf = shape.text_frame
                # Only adjust shapes that have content
                if not tf.text.strip():
                    continue

                if strategy == "shrink_text":
                    if tf.auto_size != MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
                        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        adjusted += 1
                elif strategy == "grow":
                    if tf.auto_size != MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT:
                        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                        adjusted += 1

        prs.save(file_path)

        return FixResult(
            tool_name="resize_pptx_text_boxes",
            job_id=job_id,
            success=True,
            description=f"Applied '{strategy}' to {adjusted} text box(es)",
            after_value=f"{strategy}: {adjusted} adjusted",
        )
    except Exception as exc:
        logger.error(f"Job {job_id}: resize_pptx_text_boxes failed: {exc}")
        return FixResult(
            tool_name="resize_pptx_text_boxes",
            job_id=job_id,
            success=False,
            description=f"Failed to resize PPTX text boxes: {exc}",
            error=str(exc),
        )
